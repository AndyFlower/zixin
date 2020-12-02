# -*- coding: utf-8 -*-
"""
Created on Wed Dec  2 14:22:44 2020

@author: sangliping
"""

from sys import argv
from os import listdir
from os.path import join,isfile,isdir
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
def checkdocx(dstStr,fn):
    # 打开.docx文档
    document = Document(fn)
    # 遍历所有段落文本
    for p in document.paragraphs:
        if dstStr in p.txt:
            return Ture
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                if dstStr in cell.text:
                    return True
    return False

def checkxlsx(dstStr,fn):
    wb = load_workbook(fn)
    for ws in wb.worksheets:
        for row in ws.rows:
            for cell in row:
                try:
                    if dstStr in cell.value:
                        return True
                except:
                    pass
    return False

def checkpptx(dstStr,fn):
    presentation = Presentation(fn)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.shape_type==19:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if dstStr in cell.text_frame.text:
                            return True
            elif shape.shape_type==14:
                try:
                    if dstStr in shape.text:
                        return True
                except:
                   pass
    return False

def main(dstStr,flag):
    dirs = ['.']
    while dirs:
        currentDir = dirs.pop(0)
        for fn in listdir(currentDir):
            path = join(currentDir,fn)
            if isfile(path):
                if path.endswith('.docx') and checkdocx(dstStr, path):
                    print(path)
                elif path.endswith('.xlsx') and checkxlsx(dstStr, path):
                    print(path)
                elif  path.endswith('.pptx') and checkpptx(dstStr, path):
                    print(path)
            elif flag and isdir(path):
                dirs.append(path)
                
if argv[1] !='/s':
    dstStr = argv[1]
    flag = False  
else:
    dstStr = argv[2]
    flag = True              
    
main(dstStr,flag)